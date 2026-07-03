from __future__ import annotations

from pathlib import Path


def parse_pptx(path: Path, max_slides: int = 20) -> str:
    try:
        from pptx import Presentation
    except Exception:
        return "[PPTX parser unavailable: install python-pptx]"
    prs = Presentation(str(path))
    parts = []
    for slide in prs.slides[:max_slides]:
        texts = [shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text]
        parts.append("\n".join(texts))
    return "\n\n".join(parts)

