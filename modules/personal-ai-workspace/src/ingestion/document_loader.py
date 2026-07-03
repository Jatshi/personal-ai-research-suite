from __future__ import annotations

import re
from pathlib import Path
from typing import Any

from bs4 import BeautifulSoup


def load_document(path: str | Path) -> dict[str, Any]:
    p = Path(path)
    suffix = p.suffix.lower()
    text = ""
    metadata: dict[str, Any] = {"file_name": p.name, "file_type": suffix.lstrip(".")}
    if suffix in {".md", ".txt"}:
        text = p.read_text(encoding="utf-8", errors="ignore")
    elif suffix in {".html", ".htm"}:
        soup = BeautifulSoup(p.read_text(encoding="utf-8", errors="ignore"), "html.parser")
        title = soup.title.string.strip() if soup.title and soup.title.string else p.stem
        for tag in soup(["script", "style", "nav", "footer"]):
            tag.decompose()
        text = soup.get_text("\n")
        metadata["title"] = title
    elif suffix == ".pdf":
        try:
            from pypdf import PdfReader

            reader = PdfReader(str(p))
            pages = [page.extract_text() or "" for page in reader.pages]
            text = "\n\n".join(f"[Page {i + 1}]\n{page}" for i, page in enumerate(pages))
        except Exception as exc:
            text = f"PDF parse unavailable for {p.name}: {exc}"
    elif suffix == ".docx":
        try:
            import docx

            doc = docx.Document(str(p))
            text = "\n".join(para.text for para in doc.paragraphs)
        except Exception as exc:
            text = f"DOCX parse unavailable for {p.name}: {exc}"
    elif suffix == ".pptx":
        try:
            from pptx import Presentation

            prs = Presentation(str(p))
            parts = []
            for i, slide in enumerate(prs.slides, start=1):
                parts.append(f"[Slide {i}]")
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        parts.append(shape.text)
            text = "\n".join(parts)
        except Exception as exc:
            text = f"PPTX parse unavailable for {p.name}: {exc}"
    else:
        text = p.read_text(encoding="utf-8", errors="ignore")
    metadata.setdefault("title", extract_title(text) or p.stem)
    metadata["language"] = "zh" if re.search(r"[\u4e00-\u9fff]", text) else "en"
    return {"text": text, "metadata": metadata}


def extract_title(text: str) -> str | None:
    for line in (text or "").splitlines():
        stripped = line.strip().lstrip("#").strip()
        if stripped:
            return stripped[:120]
    return None

