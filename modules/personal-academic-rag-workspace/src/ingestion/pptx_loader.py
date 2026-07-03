from __future__ import annotations

from pathlib import Path

from src.models import DocumentSegment
from src.utils.text_utils import detect_language, normalize_text


def load_pptx(path: Path, base_metadata: dict) -> list[DocumentSegment]:
    try:
        from pptx import Presentation
    except Exception as exc:
        raise RuntimeError("PPTX support requires python-pptx. Install requirements.txt.") from exc
    prs = Presentation(str(path))
    segments: list[DocumentSegment] = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text)
        text = normalize_text("\n".join(texts))
        if text:
            segments.append(
                DocumentSegment(
                    text=text,
                    metadata={**base_metadata, "page": i + 1, "paragraph": None, "language": detect_language(text)},
                )
            )
    return segments

